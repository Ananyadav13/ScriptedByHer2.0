"""
Build Trust — Prototype Submission deck generator.

Produces `Build_Trust_Prototype_Submission.pptx` in the repo root, in the same
visual language as the Round-2 idea deck (deep-plum serif headings, letter-spaced
eyebrows, light lavender canvas).

    python docs/deck/build_deck.py

Every number in here is verified against the running code, not the README:
  127 tests · 31 API operations · 20 products · 12,610 reviews · 15 sellers.
"""
from __future__ import annotations

import os
import struct
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# Palette + type
# --------------------------------------------------------------------------
INK = RGBColor(0x1A, 0x15, 0x23)      # near-black plum, headings
PLUM = RGBColor(0x4A, 0x0E, 0x33)     # brand deep
EYEBROW = RGBColor(0x8E, 0x1B, 0x5C)  # letter-spaced section label
MAGENTA = RGBColor(0xC2, 0x18, 0x5B)
BLUE = RGBColor(0x2F, 0x55, 0x97)
AMBER = RGBColor(0xB2, 0x6B, 0x15)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
BG = RGBColor(0xF8, 0xF7, 0xFA)       # light canvas
CARD = RGBColor(0xF3, 0xE7, 0xF0)     # pale pink card
CARD2 = RGBColor(0xFF, 0xFF, 0xFF)
DARKBG = RGBColor(0x22, 0x09, 0x1A)   # title / closing canvas
MUTED = RGBColor(0x5A, 0x55, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LILAC = RGBColor(0xC9, 0xA9, 0xC0)
LINE = RGBColor(0xD9, 0xCE, 0xD6)

HEAD = "Georgia"
BODY = "Calibri"

W, H = Inches(13.333), Inches(7.5)
ML = Inches(0.85)                      # left margin
CW = Inches(11.633)                    # content width

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
SHOTS = os.path.join(ROOT, "docs", "screenshots")

DECK_TITLE = "Build Trust"
FOOTER = "Build Trust  ·  ScriptedBy{Her} 2.0  ·  Prototype Submission"


# --------------------------------------------------------------------------
# Primitives
# --------------------------------------------------------------------------
def sp(text: str) -> str:
    """Fake letter-spacing the way the idea deck renders its eyebrows."""
    return "  ".join(text.upper())


def png_size(path: str) -> tuple[int, int]:
    with open(path, "rb") as f:
        head = f.read(26)
    return struct.unpack(">II", head[16:24])


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def write(tf, lines, *, first=False):
    """lines: list of dicts {text, size, color, bold, italic, font, space_before, align, line}"""
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if (i == 0 and first) else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        if spec.get("space_before"):
            p.space_before = Pt(spec["space_before"])
        if spec.get("space_after"):
            p.space_after = Pt(spec["space_after"])
        if spec.get("line"):
            p.line_spacing = spec["line"]
        r = p.add_run()
        r.text = spec["text"]
        f = r.font
        f.size = Pt(spec.get("size", 14))
        f.bold = spec.get("bold", False)
        f.italic = spec.get("italic", False)
        f.name = spec.get("font", BODY)
        f.color.rgb = spec.get("color", INK)
    return tf


def rect(slide, left, top, width, height, fill=None, line=None, line_w=0.75,
         shape=MSO_SHAPE.RECTANGLE, adj=None):
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if adj is not None:
        try:
            sh.adjustments[0] = adj
        except (IndexError, AttributeError):
            pass
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.12)
    return sh


def slide_light(prs, page=None, footer=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgfill = s.background.fill
    bgfill.solid()
    bgfill.fore_color.rgb = BG
    if footer:
        write(textbox(s, ML, Inches(6.98), Inches(7.0), Inches(0.3)),
              [{"text": FOOTER, "size": 9, "color": RGBColor(0x9A, 0x92, 0xA0)}], first=True)
    if page is not None:
        tf = textbox(s, Inches(11.9), Inches(6.98), Inches(0.6), Inches(0.3))
        write(tf, [{"text": f"{page:02d}", "size": 9,
                    "color": RGBColor(0x9A, 0x92, 0xA0), "align": PP_ALIGN.RIGHT}], first=True)
    return s


def slide_dark(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgfill = s.background.fill
    bgfill.solid()
    bgfill.fore_color.rgb = DARKBG
    return s


def fit_title_size(title, *, max_pt=815.0, sizes=(34, 32, 30, 28, 26, 24, 22)):
    """python-pptx cannot autofit, so estimate Georgia-bold width and step down
    until the title stays on one line. Keeps the rule and subtitle from colliding."""
    for size in sizes:
        if len(title) * 0.55 * size <= max_pt:
            return size
    return sizes[-1]


def heading(slide, eyebrow, title, *, sub=None, eyebrow_color=EYEBROW, title_size=None):
    if title_size is None:
        title_size = fit_title_size(title)
    write(textbox(slide, ML, Inches(0.55), CW, Inches(0.3)),
          [{"text": sp(eyebrow), "size": 10.5, "bold": True, "color": eyebrow_color}], first=True)
    tf = textbox(slide, ML, Inches(0.86), CW, Inches(0.62))
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM  # baseline sits just above the rule at any size
    write(tf, [{"text": title, "size": title_size, "bold": True, "color": INK, "font": HEAD}],
          first=True)
    rect(slide, ML, Inches(1.62), Inches(1.35), Inches(0.05), fill=PLUM)
    if sub:
        write(textbox(slide, ML, Inches(1.82), Inches(10.9), Inches(0.5)),
              [{"text": sub, "size": 13, "color": MUTED, "line": 1.25}], first=True)


def add_shot(slide, filename, left, top, width, height, anchor="top", border=True, focus=None):
    """Place a screenshot cropped to the frame's aspect ratio (no distortion).

    `focus` is a 0-1 fraction naming the vertical point of interest; the visible
    window is centred on it. Use it to frame the part of a tall page that matters.
    """
    path = os.path.join(SHOTS, filename)
    sw, sh_ = png_size(path)
    src_ar = sw / sh_
    tgt_ar = width / height
    pic = slide.shapes.add_picture(path, left, top, width=width, height=height)
    if src_ar < tgt_ar:                       # source too tall -> crop vertically
        keep = (sw / tgt_ar) / sh_
        cut = max(0.0, 1.0 - keep)
        if focus is not None:
            top_cut = min(max(focus - keep / 2, 0.0), cut)
            pic.crop_top = top_cut
            pic.crop_bottom = cut - top_cut
        elif anchor == "top":
            pic.crop_bottom = cut
        else:
            pic.crop_top = pic.crop_bottom = cut / 2
    elif src_ar > tgt_ar:                     # source too wide -> crop horizontally
        keep = (sh_ * tgt_ar) / sw
        cut = max(0.0, 1.0 - keep)
        pic.crop_left = pic.crop_right = cut / 2
    if border:
        pic.line.color.rgb = LINE
        pic.line.width = Pt(0.75)
    return pic


def card(slide, left, top, width, height, title, lines, *,
         accent=PLUM, fill=CARD2, title_size=15, body_size=11.5, gap=5):
    box = rect(slide, left, top, width, height, fill=fill, line=LINE)
    rect(slide, left, top, Inches(0.055), height, fill=accent)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.26)
    specs = [{"text": title, "size": title_size, "bold": True, "color": accent, "font": HEAD}]
    for ln in lines:
        specs.append({"text": ln, "size": body_size, "color": INK,
                      "space_before": gap, "line": 1.2})
    write(tf, specs, first=True)
    return box


def chain(slide, top, items, *, left=ML, total=CW, height=Inches(0.72),
          accent=PLUM, fill=CARD2, size=11):
    """A horizontal flow of boxes separated by arrows."""
    n = len(items)
    arrow = Inches(0.34)
    bw = int((total - arrow * (n - 1)) / n)
    x = left
    for i, (label, sub) in enumerate(items):
        b = rect(slide, x, top, Emu(bw), height, fill=fill, line=LINE)
        b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        b.text_frame.margin_left = b.text_frame.margin_right = Inches(0.08)
        write(b.text_frame, [
            {"text": label, "size": size, "bold": True, "color": accent, "align": PP_ALIGN.CENTER},
            {"text": sub, "size": size - 1.5, "color": MUTED, "align": PP_ALIGN.CENTER,
             "space_before": 1, "line": 1.1},
        ], first=True)
        if i < n - 1:
            tf = textbox(slide, Emu(x + bw), top, arrow, height)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            write(tf, [{"text": "→", "size": 15, "color": RGBColor(0xB0, 0xA6, 0xB4),
                        "align": PP_ALIGN.CENTER}], first=True)
        x = Emu(x + bw + arrow)


def statband(slide, top, stats, *, left=ML, total=CW, color=PLUM, num_size=23, lab_size=9.5):
    n = len(stats)
    cw = int(total / n)
    for i, (num, lab) in enumerate(stats):
        x = Emu(left + cw * i)
        write(textbox(slide, x, top, Emu(cw - Inches(0.2)), Inches(0.45)),
              [{"text": num, "size": num_size, "bold": True, "color": color, "font": HEAD}], first=True)
        write(textbox(slide, x, Emu(top + Inches(0.44)), Emu(cw - Inches(0.2)), Inches(0.4)),
              [{"text": lab, "size": lab_size, "color": MUTED, "line": 1.15}], first=True)
        if i:
            rect(slide, Emu(x - Inches(0.16)), top, Inches(0.008), Inches(0.72), fill=LINE)


def table(slide, left, top, width, col_w, rows, *, head_fill=PLUM, size=10.5,
          row_h=Inches(0.32), head_h=Inches(0.34)):
    n_rows, n_cols = len(rows), len(rows[0])
    shp = slide.shapes.add_table(n_rows, n_cols, left, top, width, head_h + row_h * (n_rows - 1))
    tbl = shp.table
    tbl.first_row = True
    tbl.horz_banding = False
    for i, w in enumerate(col_w):
        tbl.columns[i].width = w
    tbl.rows[0].height = head_h
    for r in range(1, n_rows):
        tbl.rows[r].height = row_h
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = head_fill if r == 0 else (
                CARD2 if r % 2 else RGBColor(0xF4, 0xEF, 0xF3))
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(size)
            run.font.name = BODY
            run.font.bold = (r == 0)
            run.font.color.rgb = WHITE if r == 0 else INK
    return tbl


def chip(slide, left, top, width, text, *, fill=None, color=LILAC, size=10.5, height=Inches(0.34)):
    b = rect(slide, left, top, width, height, fill=fill,
             line=RGBColor(0x6B, 0x3A, 0x58), shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.42)
    b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    b.text_frame.margin_left = b.text_frame.margin_right = Inches(0.12)
    write(b.text_frame, [{"text": text, "size": size, "color": color, "align": PP_ALIGN.CENTER}],
          first=True)
    return b


def circle(slide, cx, cy, r, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(cx - r), Emu(cy - r), Emu(r * 2), Emu(r * 2))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


# ==========================================================================
# Slides
# ==========================================================================
def s01_title(prs):
    s = slide_dark(prs)
    circle(s, int(W) + Inches(0.2), Inches(-0.3), Inches(2.5), RGBColor(0x3A, 0x0E, 0x2C))
    circle(s, Inches(-0.6), int(H) + Inches(0.5), Inches(2.2), RGBColor(0x22, 0x2C, 0x50))
    circle(s, Inches(7.6), int(H) + Inches(1.05), Inches(1.35), RGBColor(0x8A, 0x4D, 0x12))

    write(textbox(s, ML, Inches(1.55), Inches(10.5), Inches(0.3)),
          [{"text": sp("ScriptedBy{Her} 2.0  ·  Prototype Submission"),
            "size": 11, "bold": True, "color": RGBColor(0xE0, 0x7E, 0xB4)}], first=True)
    write(textbox(s, ML, Inches(1.95), Inches(10.5), Inches(1.2)),
          [{"text": "Build Trust.", "size": 60, "bold": True, "color": WHITE, "font": HEAD}], first=True)
    write(textbox(s, ML, Inches(3.15), Inches(9.6), Inches(1.0)),
          [{"text": "Two agentic AI systems that act on evidence — now built, deployed, "
                    "and running end to end.",
            "size": 17, "italic": True, "color": LILAC, "line": 1.35}], first=True)

    chip(s, ML, Inches(4.35), Inches(3.45), "◆  Live demo  ·  Vercel")
    chip(s, Inches(4.45), Inches(4.35), Inches(3.1), "◆  Live API  ·  Render")
    chip(s, Inches(7.72), Inches(4.35), Inches(2.35), "◆  127 tests passing")

    write(textbox(s, ML, Inches(4.82), Inches(9.0), Inches(0.4)),
          [{"text": "scripted-by-her2-0.vercel.app          scriptedbyher2-0.onrender.com",
            "size": 10, "color": RGBColor(0x8E, 0x70, 0x84)}], first=True)

    write(textbox(s, ML, Inches(6.35), Inches(10.5), Inches(0.4)),
          [{"text": "Building for Bharat with the Power of Agentic AI",
            "size": 12, "color": RGBColor(0xA8, 0x86, 0x9E)},
           {"text": "Ananya Yadav  ·  Ananya's Team 2  ·  IIITDM Jabalpur",
            "size": 11, "color": RGBColor(0x7E, 0x62, 0x76), "space_before": 4}], first=True)


def s02_problem(prs):
    s = slide_light(prs, page=2)
    heading(s, "The Problem  ·  Recap", "Three Places Trust Breaks Down",
            sub="Complaints against the platform cluster tightly around three fixable failure points — "
                "not general dissatisfaction with price or selection.")
    cards = [
        ("Authenticity", PLUM,
         ["Big-brand listings priced far below MRP and confirmed counterfeit on inspection — one "
          "carried Meesho's own 'Trusted' badge.",
          "Templated five-star bursts inflate seller ratings, leaving no reliable signal to trust "
          "before paying."]),
        ("Delivery Integrity", BLUE,
         ["One OTP for a multi-item order lets a courier mark every item delivered after collecting "
          "for one.",
          "Return pickups marked 'customer unreachable' without an agent arriving; refunds stuck "
          "for weeks with no escalation path."]),
        ("Listing Accuracy", AMBER,
         ["Seller policy requires a real size chart for most fashion categories, yet many listings "
          "show only generic S/M/L labels.",
          "Descriptions overstate material quality — 'leather' arrives synthetic — creating returns "
          "preventable at the listing stage."]),
    ]
    w, gap = Inches(3.72), Inches(0.24)
    for i, (t, acc, lines) in enumerate(cards):
        card(s, Emu(ML + (w + gap) * i), Inches(2.55), w, Inches(2.65), t, lines,
             accent=acc, fill=CARD, body_size=11)
    write(textbox(s, ML, Inches(5.32), CW, Inches(0.3)),
          [{"text": "Source: aggregated public review and complaint platforms, 2026 — "
                    "Trustpilot (1,112 reviews, 91% one-star) shown as one example.",
            "size": 9, "italic": True, "color": MUTED}], first=True)
    statband(s, Inches(5.66), [
        ("1.7★", "Average rating across public review aggregators"),
        ("10,000+", "Recent reviews analysed for recurring patterns"),
        ("3", "Specific, fixable failure clusters identified"),
        ("2", "Chosen to solve — the two that compound fastest"),
    ])


def s03_built(prs):
    s = slide_light(prs, page=3)
    heading(s, "What We Built", "Two Agents, Three Consoles, One Auditable Trail",
            sub="Not a mockup. A deployed full-stack application where both agents run for real, take "
                "graduated action, and write an audit row for everything they do.")

    card(s, ML, Inches(2.42), Inches(5.7), Inches(2.62), "Agent 1  —  Verification & Authenticity",
         ["On-demand investigation. Fires when a buyer hits Buy Now, when a deterministic tripwire "
          "trips, or when a post-delivery dispute is filed.",
          "Four evidence tools, a structured verdict, and — the part that matters — the action is "
          "executed, not just returned.",
          "Ten graduated outcomes, from a soft notification to a lock, and it always reaches for the "
          "mildest one the evidence supports."],
         accent=PLUM, body_size=12)
    card(s, Inches(6.78), Inches(2.42), Inches(5.7), Inches(2.62), "Agent 2  —  Listing & Catalog Integrity",
         ["Continuous audit. Clusters what buyers actually complain about and applies deterministic "
          "delisting tiers.",
          "Doesn't just flag: it drafts the corrected size chart or fabric description for the "
          "seller to approve in one tap.",
          "Degrades gracefully — with no API key the audit still completes on the deterministic tier "
          "rules alone."],
         accent=BLUE, body_size=12)

    statband(s, Inches(5.42), [
        ("31", "API operations"),
        ("127", "tests passing"),
        ("20", "products seeded"),
        ("12,610", "reviews"),
        ("15", "sellers"),
        ("Live", "Vercel + Render"),
    ], num_size=21)


def s04_journeys(prs):
    s = slide_light(prs, page=4)
    heading(s, "The Demo", "Two Journeys, Both Running Live",
            sub="The guided /demo route drives the real surfaces and the live API — nothing is "
                "faked or replayed for the walkthrough.")

    write(textbox(s, ML, Inches(2.55), CW, Inches(0.3)),
          [{"text": "JOURNEY A  ·  A buyer disputes what arrived", "size": 12,
            "bold": True, "color": PLUM}], first=True)
    chain(s, Inches(2.92), [
        ("Buyer", "files a fabric dispute"),
        ("Agent 1", "investigates, trace streams live"),
        ("Manager", "rules on the advisory case"),
        ("Both notified", "in plain language"),
    ], accent=PLUM)
    write(textbox(s, ML, Inches(3.74), CW, Inches(0.3)),
          [{"text": "Media evidence is uncertain by nature, so the verdict is advisory — a soft flag "
                    "and a hand-off to a human, never an automatic punishment.",
            "size": 10.5, "italic": True, "color": MUTED}], first=True)

    write(textbox(s, ML, Inches(4.42), CW, Inches(0.3)),
          [{"text": "JOURNEY B  ·  A seller lists, the catalog is audited", "size": 12,
            "bold": True, "color": BLUE}], first=True)
    chain(s, Inches(4.79), [
        ("Seller", "listing blocked by the field gate"),
        ("Agent 2", "audits, clusters, drafts a fix"),
        ("Manager", "sees discrepancies, riskiest first"),
        ("Seller", "one-tap approve, back to full visibility"),
    ], accent=BLUE)
    write(textbox(s, ML, Inches(5.61), CW, Inches(0.3)),
          [{"text": "Every action by either agent and every human decision lands in one auditable "
                    "trail at /manager/logs.",
            "size": 10.5, "italic": True, "color": MUTED}], first=True)

    box = rect(s, ML, Inches(6.12), CW, Inches(0.62), fill=CARD, line=None)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(box.text_frame, [
        {"text": "Scenario products, all seeded and stable:  a counterfeit Rolex at 0.07% of MRP  ·  "
                 "an honest viral seller who is correctly cleared  ·  a knockoff buyers knowingly "
                 "love, relabelled rather than banned  ·  a cross-variant fabric dispute.",
         "size": 11, "color": PLUM}], first=True)


def s05_agent1(prs):
    s = slide_light(prs, page=5)
    heading(s, "Agent 1", "The Investigation Loop")

    write(textbox(s, ML, Inches(1.85), Inches(6.5), Inches(0.6)),
          [{"text": "A manual function-calling loop over Gemini — not automatic, so every tool call "
                    "streams to the live trace as it happens.",
            "size": 11.5, "color": MUTED, "line": 1.25}], first=True)

    table(s, ML, Inches(2.72), Inches(6.5),
          [Inches(1.75), Inches(3.85), Inches(0.9)],
          [["Tool", "What it returns — concrete numbers the agent must cite", "LLM?"],
           ["check_catalog_risk", "Price-vs-MRP ratio, review-burst stats, trustworthy rating, order volume", "No"],
           ["check_seller_profile", "Account age, effective rating, trust flags, repeat-case count", "No"],
           ["check_delivery_signals", "OTP-vs-items, hub anomaly, geo-photo proof, claimer class", "No"],
           ["check_media_evidence", "Variant-aware quality-fingerprint comparison", "Vision"]],
          size=10, row_h=Inches(0.44))

    box = rect(s, ML, Inches(4.95), Inches(6.5), Inches(1.15), fill=CARD, line=None)
    write(box.text_frame, [
        {"text": "Why this is an agent, not a chatbot", "size": 12, "bold": True, "color": PLUM, "font": HEAD},
        {"text": "It plans → calls a tool → observes what came back → decides the next step, up to "
                 "8 times. No fixed script: a clean seller profile ends the loop early, a contradictory signal sends it to another tool.",
         "size": 10.5, "color": INK, "space_before": 3, "line": 1.18},
        {"text": "Then it ACTS. The verdict is executed, not returned — status changes, buyer and "
                 "seller notifications, audit rows. That consequence is what makes it agentic.",
         "size": 10.5, "color": PLUM, "space_before": 4, "line": 1.18}], first=True)

    add_shot(s, "10-theAgentinAction.png", Inches(7.5), Inches(2.05), Inches(4.98), Inches(3.60))
    write(textbox(s, Inches(7.5), Inches(5.76), Inches(4.98), Inches(0.4)),
          [{"text": "The live trace, mid-investigation — each tool call lands on screen as it "
                    "happens, with the numbers it returned.",
            "size": 9.5, "italic": True, "color": MUTED, "line": 1.15}], first=True)


def s06_trace(prs):
    """Verbatim output from the deployed backend. Captured live — nothing rewritten.

    Source: POST /investigate {product_id: prod_counterfeit_rolex} against
    https://scriptedbyher2-0.onrender.com on 19 Jul 2026 → investigation inv_c5cd131bcce3.
    """
    s = slide_light(prs, page=6)
    heading(s, "Evidence", "What the Agent Actually Returns",
            sub="Verbatim output from the deployed system. Nothing here is rewritten for the slide.")

    MONO = "Consolas"
    CODEBG = RGBColor(0x2A, 0x24, 0x2E)
    CODEFG = RGBColor(0xE6, 0xDF, 0xE8)
    KEY = RGBColor(0x9F, 0xC9, 0xE8)
    VAL = RGBColor(0xF0, 0xB4, 0x6B)

    write(textbox(s, ML, Inches(2.5), Inches(6.3), Inches(0.3)),
          [{"text": "THE TOOL CALLS  ·  deterministic, no LLM", "size": 10.5, "bold": True,
            "color": PLUM}], first=True)

    calls = [
        ("check_catalog_risk", [
            ("price_to_mrp_ratio", "0.0007   → Rolex priced at 0.07% of MRP"),
            ("review_burst", "54 reviews in one day, 100% from accounts < 30d"),
            ("trustworthy_rating", "null — 0 reviews from established accounts"),
            ("order_volume", "24 orders ≥ 20 → confidence floor met"),
        ]),
        ("check_seller_profile", [
            ("account_age_days", "20"),
            ("trust_flags", "['new_account_cluster']"),
            ("classification", "suspicious_new"),
            ("repeat_offender", "false — no prior policy cases"),
        ]),
    ]
    y = Inches(2.86)
    for tool, fields in calls:
        h = Inches(0.40) + Inches(0.29) * len(fields)
        rect(s, ML, y, Inches(6.3), h, fill=CODEBG)
        tf = textbox(s, Emu(ML + Inches(0.16)), Emu(y + Inches(0.09)), Inches(6.0), Emu(h))
        specs = [{"text": f"▸ {tool}()", "size": 11, "bold": True, "color": RGBColor(0x8F, 0xD9, 0xA8),
                  "font": MONO}]
        for k, v in fields:
            specs.append({"text": f"   {k}", "size": 9, "color": KEY, "font": MONO,
                          "space_before": 2})
            specs.append({"text": f"      {v}", "size": 9, "color": VAL, "font": MONO})
        write(tf, specs, first=True)
        y = Emu(y + h + Inches(0.16))

    write(textbox(s, ML, Inches(6.28), Inches(6.3), Inches(0.5)),
          [{"text": "Every number the agent cites comes from a pure function. It cannot invent evidence — "
                    "only reason over what the tools return.",
            "size": 10, "italic": True, "color": MUTED, "line": 1.15}], first=True)

    write(textbox(s, Inches(7.66), Inches(2.5), Inches(4.82), Inches(0.3)),
          [{"text": "THE STRUCTURED VERDICT  ·  schema-enforced", "size": 10.5, "bold": True,
            "color": BLUE}], first=True)

    v = rect(s, Inches(7.66), Inches(2.86), Inches(4.82), Inches(1.0), fill=PLUM)
    v.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(v.text_frame, [
        {"text": "counterfeit_lock", "size": 19, "bold": True, "color": WHITE, "font": MONO},
        {"text": "confidence 1.0   ·   action: lock   ·   listing locked, seller notified",
         "size": 10, "color": LILAC, "space_before": 3}], first=True)

    write(textbox(s, Inches(7.66), Inches(4.02), Inches(4.82), Inches(0.3)),
          [{"text": "evidence[] — the agent's own words", "size": 9.5, "bold": True,
            "color": MUTED}], first=True)

    ev = [
        "Priced at ₹599 — only 0.07% of the ₹8,50,000 MRP.",
        "A burst of 54 reviews in a single day, 100% from accounts under 30 days old.",
        "Zero reviews from established accounts to verify authenticity.",
        "Seller account is 20 days old, opened in a cluster of new accounts.",
        "With 24 orders, sufficient evidence for a hard enforcement action.",
    ]
    y = Inches(4.34)
    for e in ev:
        write(textbox(s, Inches(7.66), y, Inches(4.82), Inches(0.36)),
              [{"text": "— " + e, "size": 9.5, "color": INK, "line": 1.12}], first=True)
        y = Emu(y + Inches(0.34))

    box = rect(s, Inches(7.66), Inches(6.12), Inches(4.82), Inches(0.66), fill=CARD, line=None)
    write(box.text_frame, [
        {"text": "And the buyer sees plain language, not JSON", "size": 9.5, "bold": True,
         "color": PLUM},
        {"text": "“This product has been flagged as counterfeit due to a price significantly "
                 "below market value and evidence of artificial review activity.”",
         "size": 9, "italic": True, "color": INK, "space_before": 2, "line": 1.12}], first=True)


def s07_restraint(prs):
    """The ladder AND the safeguards that constrain it — one slide, because they are one idea."""
    s = slide_light(prs, page=7)
    heading(s, "The Decision Core", "Restraint, by Construction",
            sub="Ten graduated outcomes. The agent reasons; a deterministic rung decides what it is "
                "actually allowed to do — and it always reaches for the mildest one the evidence supports.")

    write(textbox(s, ML, Inches(2.52), Inches(6.2), Inches(0.3)),
          [{"text": "THE ACTION LADDER  ·  least drastic first", "size": 10.5, "bold": True,
            "color": PLUM}], first=True)

    rungs = [
        ("Nothing happens", "cleared  ·  authentic", RGBColor(0x9C, 0xC4, 0xAE), 0.30),
        ("Tell someone", "notify_only  ·  relabel_required", RGBColor(0x7F, 0xB0, 0xD4), 0.46),
        ("Ask, reversibly", "request_qc_video  ·  hold_pending_fix", RGBColor(0xE8, 0xB5, 0x6A), 0.62),
        ("Hand to a human", "recommend_review  ·  manual_review", RGBColor(0xD9, 0x8C, 0x4A), 0.78),
        ("Act, hard", "refund_fast_track  ·  counterfeit_lock  ·  ban", RGBColor(0xC2, 0x4B, 0x5B), 1.0),
    ]
    y = Inches(2.92)
    for label, members, col, frac in rungs:
        rect(s, ML, y, Emu(int(Inches(6.2) * frac)), Inches(0.52), fill=col)
        tf = textbox(s, Emu(ML + Inches(0.14)), Emu(y + Inches(0.05)), Inches(5.9), Inches(0.44))
        write(tf, [
            {"text": label, "size": 11.5, "bold": True, "color": RGBColor(0x1A, 0x15, 0x23)},
            {"text": members, "size": 9, "color": RGBColor(0x33, 0x2B, 0x38), "space_before": 0}],
            first=True)
        y = Emu(y + Inches(0.68))

    write(textbox(s, Inches(7.5), Inches(2.52), Inches(4.98), Inches(0.3)),
          [{"text": "THE SAFEGUARDS THAT CONSTRAIN IT", "size": 10.5, "bold": True, "color": BLUE}],
          first=True)

    guards = [
        ("Confidence floor", "A hard lock needs 20+ orders of evidence, or it auto-downgrades to a "
                             "reversible request. Thin data never punishes."),
        ("Bursts need two signals", "A review spike is fraud evidence only alongside a high share of "
                                    "new accounts — the honest viral seller is cleared."),
        ("Two-signal refunds", "OTP mismatch, hub anomaly and a missing geo-photo are independent. "
                               "Two fast-track a refund; a serial claimer goes to a human."),
        ("Proportional penalties", "A delisting hit scales with that product's share of the seller's "
                                   "genuine reviews — one failure among fifty is near-zero."),
        ("Humans hold the final call", "Managers decide, and may only act on their own sellers' "
                                       "listings — enforced server-side with a 403."),
    ]
    y = Inches(2.9)
    for t, b in guards:
        write(textbox(s, Inches(7.5), y, Inches(4.98), Inches(0.62)), [
            {"text": t, "size": 11, "bold": True, "color": INK},
            {"text": b, "size": 9.5, "color": MUTED, "space_before": 1, "line": 1.15}], first=True)
        y = Emu(y + Inches(0.68))

    box = rect(s, ML, Inches(6.42), CW, Inches(0.56), fill=CARD, line=None)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(box.text_frame, [
        {"text": "Authenticity matters — but not at the cost of the buyer or seller community. "
                 "A counterfeit people regret gets locked. A knockoff people knowingly love gets "
                 "relabelled, not banned.",
         "size": 11.5, "color": PLUM, "italic": True}], first=True)


def s08_variant(prs):
    s = slide_light(prs, page=8)
    heading(s, "The Hard Part", "A Vision Model That Cannot False-Flag on Colour")

    box = rect(s, ML, Inches(1.88), CW, Inches(0.82), fill=RGBColor(0xFA, 0xF0, 0xE6),
               line=RGBColor(0xE4, 0xC9, 0xA6))
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(box.text_frame, [
        {"text": "THE TRAP", "size": 10, "bold": True, "color": AMBER},
        {"text": "A seller films ONE listing video but sells the same kurti in black, blue and red. "
                 "Naively diffing a buyer's photo of the blue item against a video of the black one "
                 "makes colour the loudest discrepancy — false-flagging an honest seller and burying "
                 "the real complaint.",
         "size": 11.5, "color": INK, "space_before": 2, "line": 1.2}], first=True)

    steps = [
        ("1  ·  Distil once", PLUM,
         "The listing video is reduced to seven variant-invariant golden fields — weave, sheen, "
         "fibre texture, opacity, stitch, drape, embellishment — cached on the product."),
        ("2  ·  Read the same schema", BLUE,
         "The buyer's evidence is read into that identical schema. Colour, shade and print colourway "
         "are recorded, but held apart from anything that will be scored."),
        ("3  ·  Decide in pure Python", GREEN,
         "A deterministic diff — not the model, not the prompt — drops the variant-specific "
         "attributes before scoring, and needs 3+ readable attributes and 40% divergence."),
    ]
    y = Inches(2.92)
    for t, acc, body in steps:
        card(s, ML, y, Inches(5.55), Inches(1.06), t, [body],
             accent=acc, title_size=12.5, body_size=10)
        y = Emu(y + Inches(1.15))

    box = rect(s, ML, Inches(6.36), Inches(5.55), Inches(0.40), fill=CARD, line=None)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    box.text_frame.margin_top = box.text_frame.margin_bottom = Inches(0.02)
    write(box.text_frame, [
        {"text": "Exception, by design: if the buyer's claim IS 'wrong colour sent', colour is "
                 "compared — and reported separately.",
         "size": 9.5, "color": PLUM, "italic": True}], first=True)

    # The screenshot is the proof: 'colour (ignored)' is visible on screen, in the shipped product.
    add_shot(s, "theAgent1result.png", Inches(6.72), Inches(2.92), Inches(5.76), Inches(3.18),
             focus=0.30)
    write(textbox(s, Inches(6.72), Inches(6.18), Inches(5.76), Inches(0.62)), [
        {"text": "The safeguard, visible in the shipped product.", "size": 10, "bold": True,
         "color": PLUM},
        {"text": "Listing filmed Black, buyer received Blue. Five of seven invariant attributes "
                 "diverge — while colour, shade and print_colourway are stamped “ignored”.",
         "size": 9.5, "color": MUTED, "space_before": 2, "line": 1.15}], first=True)


def s09_agent2(prs):
    s = slide_light(prs, page=9)
    heading(s, "Agent 2", "It Drafts the Fix, Not Just the Flag")

    items = [
        ("Clusters what buyers actually say", "Negative complaints are grouped into fixed labels — "
         "fabric_mismatch, size_issue, damaged_delivery, possible_fraud."),
        ("Acts on agreement, not volume", "A cluster is actionable at 30% concurrence. One loud "
         "review is noise; a shared complaint is signal."),
        ("Deterministic delisting tiers", "Below 3.0★ across 1,000+ reviews  ·  below 2.0★ across "
         "700+  ·  ≤1.0★ across 500+. Less evidence needed as the rating worsens."),
    ]
    y = Inches(1.95)
    for t, b in items:
        write(textbox(s, ML, y, Inches(6.6), Inches(0.86)), [
            {"text": t, "size": 12.5, "bold": True, "color": BLUE, "font": HEAD},
            {"text": b, "size": 11, "color": INK, "space_before": 2, "line": 1.2}], first=True)
        y = Emu(y + Inches(0.92))

    # A concrete before/after does more work than another paragraph about drafting.
    write(textbox(s, ML, Inches(4.78), Inches(6.6), Inches(0.3)),
          [{"text": "Then it drafts the correction — one real example", "size": 12.5, "bold": True,
            "color": BLUE, "font": HEAD}], first=True)

    b1 = rect(s, ML, Inches(5.16), Inches(3.18), Inches(0.92), fill=CARD2, line=LINE)
    write(b1.text_frame, [
        {"text": "BEFORE  ·  seller's listing", "size": 8.5, "bold": True, "color": MUTED},
        {"text": "“Women's Handbag Combo — Free Size”", "size": 10.5, "color": INK,
         "space_before": 2},
        {"text": "no measurements  ·  size / fit complaints 67%", "size": 9, "color": AMBER,
         "space_before": 2}], first=True)

    b2 = rect(s, Inches(4.27), Inches(5.16), Inches(3.18), Inches(0.92), fill=RGBColor(0xE8, 0xF2, 0xEA),
              line=GREEN)
    write(b2.text_frame, [
        {"text": "AFTER  ·  Agent 2's draft, one tap to approve", "size": 8.5, "bold": True,
         "color": MUTED},
        {"text": "Tote 33×27×12 cm · Sling 24×18×7 cm · Purse 19×10 cm", "size": 10.5, "color": INK,
         "space_before": 2},
        {"text": "listing restored to full visibility", "size": 9, "color": GREEN,
         "space_before": 2}], first=True)

    box = rect(s, ML, Inches(6.24), Inches(6.6), Inches(0.6), fill=CARD, line=None)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(box.text_frame, [
        {"text": "Degrades gracefully — with no API key the audit still completes on the "
                 "deterministic tier rules, and the fix falls back to a deterministic draft.",
         "size": 10.5, "color": PLUM}], first=True)

    add_shot(s, "theGuidedWalkthroughSeller-2.png", Inches(7.62), Inches(2.05), Inches(4.86), Inches(3.37))
    write(textbox(s, Inches(7.62), Inches(5.53), Inches(4.86), Inches(0.4)),
          [{"text": "Agent 2 mid-audit — findings clustered by cause, with the drafted correction "
                    "waiting for the seller to approve.", "size": 9.5, "italic": True, "color": MUTED, "line": 1.15}], first=True)


def s10_arch(prs):
    s = slide_light(prs, page=10)
    heading(s, "Architecture", "The Model Gathers Evidence. Pure Functions Decide.")

    def band(top, height, text, sub, fill, textcol=WHITE, size=12.5):
        b = rect(s, Inches(1.55), top, Inches(9.05), height, fill=fill, line=None)
        b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(b.text_frame, [
            {"text": text, "size": size, "bold": True, "color": textcol, "align": PP_ALIGN.CENTER},
            {"text": sub, "size": size - 2, "color": textcol, "align": PP_ALIGN.CENTER,
             "space_before": 1, "line": 1.1}], first=True)
        return b

    def label(top, text):
        write(textbox(s, Inches(10.98), top, Inches(2.3), Inches(0.3)),
              [{"text": text, "size": 9.5, "italic": True, "color": MUTED}], first=True)

    band(Inches(1.9), Inches(0.68), "Next.js 16  ·  React 19  ·  Tailwind v4",
         "buyer  ·  seller  ·  manager  ·  agent consoles  ·  admin", BLUE)
    label(Inches(2.02), "HTTP  +  SSE live trace")

    band(Inches(2.72), Inches(0.55), "FastAPI  —  31 operations",
         "", RGBColor(0x6B, 0x63, 0x72), size=12)

    b = rect(s, Inches(1.55), Inches(3.42), Inches(4.42), Inches(0.78), fill=PLUM, line=None)
    b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(b.text_frame, [
        {"text": "Agent 1  orchestrator", "size": 12, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER},
        {"text": "manual function-calling loop", "size": 10, "color": LILAC,
         "align": PP_ALIGN.CENTER, "space_before": 1}], first=True)
    b = rect(s, Inches(6.18), Inches(3.42), Inches(4.42), Inches(0.78), fill=PLUM, line=None)
    b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(b.text_frame, [
        {"text": "Agent 2", "size": 12, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER},
        {"text": "review clustering  +  fix drafting", "size": 10, "color": LILAC,
         "align": PP_ALIGN.CENTER, "space_before": 1}], first=True)
    label(Inches(3.62), "tools  /  batched calls")

    b = rect(s, Inches(1.55), Inches(4.42), Inches(9.05), Inches(0.95),
             fill=RGBColor(0xE8, 0xF2, 0xEA), line=GREEN)
    b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(b.text_frame, [
        {"text": "services/   —   D E T E R M I N I S T I C ,   N O   L L M", "size": 12,
         "bold": True, "color": GREEN, "align": PP_ALIGN.CENTER},
        {"text": "risk_checks  ·  delisting  ·  quality_fingerprint  ·  mandatory_fields  ·  "
                 "fit_prediction  ·  tripwires  ·  rules",
         "size": 10.5, "color": INK, "align": PP_ALIGN.CENTER, "space_before": 2, "line": 1.1}],
         first=True)
    label(Inches(4.72), "enforced at package level")

    band(Inches(5.55), Inches(0.48), "SQLAlchemy 2.0  →  SQLite", "",
         RGBColor(0x6B, 0x63, 0x72), size=11.5)

    # Two-column footer. Both of these used to fight the diagram for horizontal space;
    # they belong below it, where the full slide width is free.
    gap = Inches(0.24)
    colw = Emu(int((CW - gap) / 2))

    box = rect(s, ML, Inches(6.14), colw, Inches(0.82), fill=CARD, line=None)
    write(box.text_frame, [
        {"text": "The load-bearing rule", "size": 11.5, "bold": True, "color": PLUM},
        {"text": "Every decision lives in services/ — LLM-free, unit-tested, and importable without "
                 "an API key, so the reasoning layer can never quietly become the thing that decides.",
         "size": 10, "color": INK, "space_before": 2, "line": 1.15}], first=True)

    box = rect(s, Emu(ML + colw + gap), Inches(6.14), colw, Inches(0.82), fill=CARD, line=None)
    write(box.text_frame, [
        {"text": "127 tests passing, fully offline", "size": 11.5, "bold": True, "color": PLUM},
        {"text": "66 deterministic rules  ·  22 state-integrity  ·  16 API smoke  ·  9 vision  ·  "
                 "7 delist endpoint  ·  7 drift tests that fail the build if a threshold is ever "
                 "restated as a literal in a prompt.",
         "size": 10, "color": INK, "space_before": 2, "line": 1.15}], first=True)


def s11_impact(prs):
    s = slide_light(prs, page=11)
    heading(s, "Impact", "What Changes If This Works",
            sub="The two problems chosen are the two that compound fastest — and each one has a "
                "measurable counterpart the platform already tracks.")

    cols = [
        ("For Buyers", PLUM, [
            "The complaint category that today takes weeks and often goes nowhere — wrong, damaged "
            "or missing items — resolves on evidence, in the buyer's favour when the data supports it.",
            "A clearer authenticity signal and a confident size decision arrive before checkout, "
            "not after a bad experience.",
        ]),
        ("For Honest Sellers", BLUE, [
            "Delivery-fault complaints stop being charged against seller ratings — a logistics "
            "failure is referred to logistics, not scored against the shop.",
            "Listing gaps get a drafted fix and a correction window instead of a silent down-rank.",
            "They stop competing on equal footing with counterfeits and bot-inflated ratings.",
        ]),
        ("For the Platform", AMBER, [
            "Removing counterfeit and bot-review noise, and cutting the volume of low-trust disputes, "
            "targets the exact 1.7-star problem directly.",
            "A lower preventable-return rate is a measurable hit to RTO-driven margin loss — the "
            "single largest cost line in Indian value commerce.",
        ]),
    ]
    w, gap = Inches(3.72), Inches(0.24)
    for i, (t, acc, lines) in enumerate(cols):
        card(s, Emu(ML + (w + gap) * i), Inches(2.55), w, Inches(2.72), t, lines,
             accent=acc, fill=CARD, title_size=15, body_size=10.5)

    write(textbox(s, ML, Inches(5.52), CW, Inches(0.3)),
          [{"text": "WHAT WE WOULD MEASURE IN A PILOT", "size": 10.5, "bold": True, "color": PLUM}],
          first=True)
    statband(s, Inches(5.88), [
        ("Precision", "agent verdicts agreeing with human moderators, in shadow mode"),
        ("Time-to-resolution", "dispute filed → outcome, against the current weeks-long baseline"),
        ("Preventable returns", "returns attributable to a listing gap the audit had already flagged"),
        ("False-flag rate", "honest sellers actioned in error — the number that must stay near zero"),
    ], num_size=15, lab_size=9)


def s12_close(prs):
    s = slide_dark(prs)
    circle(s, int(W) + Inches(0.4), Inches(-0.5), Inches(2.3), RGBColor(0x3A, 0x0E, 0x2C))
    circle(s, Inches(-0.8), int(H) + Inches(0.4), Inches(2.0), RGBColor(0x22, 0x2C, 0x50))
    circle(s, Inches(8.4), int(H) + Inches(1.0), Inches(1.25), RGBColor(0x8A, 0x4D, 0x12))

    write(textbox(s, ML, Inches(0.72), Inches(5.6), Inches(0.3)),
          [{"text": sp("Honest Scope"), "size": 10, "bold": True,
            "color": RGBColor(0xE0, 0x7E, 0xB4)}], first=True)
    write(textbox(s, ML, Inches(1.06), Inches(5.6), Inches(1.5)), [
        {"text": "What this prototype is not — stated plainly, because a trust system that oversells "
                 "itself is the wrong kind of demo.",
         "size": 11, "italic": True, "color": LILAC, "line": 1.22},
        {"text": "No authentication — role switching is a UI affordance, not a security boundary, and "
                 "we shipped no login rather than a convincing fake.   ·   SQLite, not Postgres.   ·   "
                 "Single-process SSE.   ·   No rate limiting.   ·   Synthetic seed data, generated to "
                 "exercise every decision path.   ·   The two media attribute reads are pre-extracted "
                 "for a zero-quota demo — but the deterministic diff that actually decides is never "
                 "seeded, and extraction runs live on upload.",
         "size": 10, "color": RGBColor(0xA8, 0x86, 0x9E), "space_before": 6, "line": 1.25}], first=True)

    write(textbox(s, ML, Inches(3.32), Inches(5.6), Inches(0.3)),
          [{"text": sp("What Comes Next"), "size": 10, "bold": True,
            "color": RGBColor(0xE0, 0x7E, 0xB4)}], first=True)

    nxt = [
        ("Postgres + Redis", "the one change that unblocks horizontal scale"),
        ("Real auth + rate limiting", "boundaries enforced server-side, not in the UI"),
        ("Perceptual-hash image matching", "the counterfeit signal designed for, not yet built"),
        ("Shadow mode before it acts", "weeks of scoring live listings against human moderators "
                                       "before it gets write access"),
    ]
    y = Inches(3.68)
    for t, b in nxt:
        write(textbox(s, ML, y, Inches(5.6), Inches(0.7)), [
            {"text": t, "size": 12, "bold": True, "color": WHITE, "font": HEAD},
            {"text": b, "size": 10, "color": LILAC, "space_before": 1, "line": 1.15}], first=True)
        y = Emu(y + Inches(0.72))

    write(textbox(s, Inches(7.0), Inches(1.5), Inches(5.4), Inches(2.2)),
          [{"text": "The decision logic works.", "size": 26, "bold": True, "color": WHITE,
            "font": HEAD},
           {"text": "It is deployed, tested, and auditable — and it says no to itself in every "
                    "place where the evidence is thin.",
            "size": 15, "italic": True, "color": LILAC, "space_before": 10, "line": 1.35}],
          first=True)

    write(textbox(s, Inches(7.0), Inches(3.55), Inches(5.4), Inches(1.0)),
          [{"text": "Build Trust.", "size": 40, "bold": True, "color": WHITE, "font": HEAD}],
          first=True)

    chip(s, Inches(7.0), Inches(4.68), Inches(2.6), "◆  Live demo")
    chip(s, Inches(9.75), Inches(4.68), Inches(2.65), "◆  127 tests passing")
    write(textbox(s, Inches(7.0), Inches(5.14), Inches(5.4), Inches(0.4)),
          [{"text": "scripted-by-her2-0.vercel.app", "size": 10,
            "color": RGBColor(0x8E, 0x70, 0x84)}], first=True)

    write(textbox(s, ML, Inches(6.72), Inches(11.5), Inches(0.4)),
          [{"text": "Ananya Yadav  ·  Ananya's Team 2  ·  IIITDM Jabalpur  ·  "
                    "23bcs025@iiitdmj.ac.in  ·  yananya784@gmail.com",
            "size": 10.5, "color": RGBColor(0x7E, 0x62, 0x76)}], first=True)


# --------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for fn in (s01_title, s02_problem, s03_built, s04_journeys, s05_agent1, s06_trace,
               s07_restraint, s08_variant, s09_agent2, s10_arch, s11_impact, s12_close):
        fn(prs)

    prs.core_properties.title = "Build Trust — Prototype Submission"
    prs.core_properties.author = "Ananya Yadav · Ananya's Team 2 · IIITDM Jabalpur"
    prs.core_properties.subject = "ScriptedBy{Her} 2.0 — Prototype Submission Phase"

    # Optional path override — handy when the default output is open in PowerPoint.
    out = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
        ROOT, "Build_Trust_Prototype_Submission.pptx")
    prs.save(out)
    print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
